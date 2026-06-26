"""将 HTML 幻灯片的设计参数迁移到 PPTX 文件。

此脚本不解析 HTML，而是直接使用 HTML 中确定的设计参数。
用户需要手动将论文截图插入到对应占位框中。
所有占位框在 PPTX 中以文本框标注来源 + PDF 位置百分比。
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ---- 设计参数（与 slides.html CSS 变量一致）----
DEEP_BLUE = RGBColor(0x1A, 0x36, 0x5D)
AMBER = RGBColor(0xA0, 0x52, 0x0A)  # darkened for contrast
GRAPHITE = RGBColor(0x37, 0x41, 0x51)
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
TITANIUM = RGBColor(0xF7, 0xF8, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = 'Source Han Sans SC'
FONT_BODY = 'Source Han Serif SC'
FONT_MONO = 'JetBrains Mono'

TITLE_SIZE = Pt(28)
BODY_SIZE = Pt(16)
SMALL_SIZE = Pt(10)
TAG_SIZE = Pt(10)

# ---- 44 页幻灯片定义 ----
SLIDES = [
    # === 开篇 (1-5) ===
    (1, "汇报主题", "黑滑石高值化商业化初步研究汇报", "cover",
     "封面页 — 深蓝渐变背景 + 汇报人/导师/日期信息",
     [("背景图", "广丰黑滑石矿区全景", "")]),

    (2, "汇报提纲", "汇报提纲", "list",
     "6 个章节目录 — 01 背景 · 02 知识资产 · 03 初阶方案 · 04 中阶展望 · 05 高阶愿景 · 06 行动讨论",
     []),

    (3, "开篇 · 背景", "黑滑石：被低估的 10 亿吨", "two-col",
     "左：原矿照片占位 | 右：矿物学身份卡片(化学式/结构/MgO含量/SiO₂含量/有机碳/储量)",
     [("黑滑石原矿手标本照片", "PAPER-005 Meng 2022 Fig.1", "PDF 第 2 页, 约 15% 位置")]),

    (4, "开篇 · 问题", "核心矛盾：为什么 10 亿吨堆在那里？", "two-col",
     "左：当前困境(陶瓷填料80%+/煅烧白滑石¥500-800/吨/黑色低30-50%) | 右：新视角(碳原位石墨烯化/Mg分离/Si陶瓷化/导电填料¥5,000-20,000/吨)",
     []),

    (5, "开篇 · 路线全景", "五条技术路线：覆盖从低到高的阶梯", "svg-chart",
     "SVG流程图 — 黑滑石原矿→D1传统煅烧/D2 5G微波陶瓷/D3 N₂碳石墨烯(初阶旗帜)/D4抗菌岩板/D9酸浸分离",
     []),

    # === 痛点即机会 (6-10) ===
    (6, "痛点即机会 · 现状", "广丰黑滑石：储量 vs 价值的巨大落差", "two-col",
     "左：矿区地图占位 + XRF成分表占位 | 右：数据卡片(储量10亿吨/MgO 23-31%/售价~200-400元/吨)",
     [("广丰矿区地理位置(Google Earth)", "Google Earth 截图", ""),
      ("XRF 成分表", "PAPER-005 Meng 2022 Table 1", "PDF 第 3 页, 约 20% 位置")]),

    (7, "痛点即机会 · 认知", "碳：从杂质到资源的认知翻转", "svg-chart",
     "分叉对比图 — 空气煅烧→CO₂→白滑石¥500-800/吨 vs N₂气氛→原位石墨烯化→导电填料¥5,000-20,000/吨 · 抗弯+43%/电阻率降1数量级",
     []),

    (8, "痛点即机会 · 科学基础", "黑滑石热演化：PAPER-005 Meng 的核心发现", "two-col",
     "左：TG-DSC + In-situ XRD 占位 | 右：四阶段SVG示意图(Stage I-IV) + 活化窗口800-900°C + N₂保护窗口>165°C",
     [("TG-DSC 热分析曲线", "PAPER-005 Meng 2022 Fig.2", "PDF 第 5 页, 约 35% 位置"),
      ("In-situ XRD 物相演化", "PAPER-005 Meng 2022 Fig.6", "PDF 第 9 页, 约 60% 位置")]),

    (9, "痛点即机会 · 验证", "N₂ vs Air：同一块矿石的两种命运", "two-col",
     "左：SVG柱状图(抗弯51.8 vs 36.1 MPa +43%/电阻率降1数量级/导热+0.62 W/(m·K)) | 右：TEM石墨烯占位 + 机理说明",
     [("N₂ vs Air 性能对比数据", "PAPER-006 吴小文 2017 Fig.3-5", "PDF 约 45-55% 位置"),
      ("TEM 石墨烯高分辨像", "PAPER-006 吴小文 2017", "PDF 约 70% 位置")]),

    (10, "痛点即机会 · 总结", "三个科学事实指向一个商业机会", "svg-chart",
     "三圆汇聚图 — 10亿吨储量 × 四阶段热演化模型 × N₂碳石墨烯化 → 低成本原位石墨烯功能填料产业化",
     []),

    # === 知识资产 (11-15) ===
    (11, "知识资产 · 论文总览", "10 篇论文：五条路线的科学地基", "cards",
     "2行×5列论文卡片矩阵 — 每卡含论文编号/作者/年份/核心贡献,按路线颜色区分(D1灰/D2蓝/D3绿/D4紫/D9橙)",
     []),

    (12, "知识资产 · 数据快照", "从论文中提取的核心数据", "table",
     "按路线(基础/D2/D3/D4/D9)分组的数据表格 — 包含Q×f/温度/成本/浸出率等关键指标",
     []),

    (13, "知识资产 · 关联网络", "25 条论文间关联：知识不是孤岛", "svg-chart",
     "SVG节点-连线图 — PAPER-002和PAPER-005为中心节点,连接D2簇(5篇)/D3/D4/D9簇",
     []),

    (14, "知识资产 · D2 深度", "D2 路线：五篇论文构成完整技术链", "svg-chart",
     "水平流程 — PAPER-004(基础验证)→PAPER-008(成本突破)→PAPER-009(低温烧结)→PAPER-010+011(掺杂优化)→Cu+Co共掺杂=D2最优解",
     []),

    (15, "知识资产 · 金字塔", "知识资产的三个层次", "svg-chart",
     "三层金字塔 — Layer1:PAPER-005四阶段模型(地基) / Layer2:PAPER-002综述框架 / Layer3:各路线实验验证(应用)",
     []),

    # === 初阶深讲 (16-33) ===
    (16, "初阶方案 · 总览", "初阶方案：现在就能启动的四个阶段", "cards",
     "四列卡片 — Phase 0a文献综述(¥0/1-2月) / 0b数据再分析(¥0/1月) / 0c管式炉实验(¥5,000/2-3月) / 0d专利申请(¥3,000/1-2月) · 总预算~8,000元/1-2学期/产出3论文+1专利",
     []),

    (17, "初阶方案 · Phase 0a", "Phase 0a：文献综述 · 站在 10 篇论文的肩膀上", "table",
     "5篇可直接引用的核心论文 — PAPER-002/PAPER-005/PAPER-006/PAPER-007/PAPER-003 与对应的综述章节和关键引用数据",
     []),

    (18, "初阶方案 · Phase 0a", "Phase 0a：系统检索策略", "list",
     "中英文检索式(WoS/CNKI/Google Scholar) + 目标期刊(Carbon/Appl Clay Sci/Ceramics Int等) + 预期: 5-15篇直接相关论文 · 产出:1篇中文综述初稿(材料导报/硅酸盐通报)",
     []),

    (19, "初阶方案 · Phase 0b", "Phase 0b：温度-气氛-产物 三位一体相图", "two-col",
     "左：3个论文截图占位(TG-DSC/In-situ XRD/性能对比) | 右：空气 vs N₂气氛四阶段对比相图 + 产物:白色惰性填料 vs 灰色导电/导热填料",
     [("TG-DSC 曲线", "PAPER-005 Meng 2022 Fig.2", "PDF 约 35% 位置"),
      ("In-situ XRD", "PAPER-005 Meng 2022 Fig.6", "PDF 约 60% 位置"),
      ("性能对比数据", "PAPER-006 吴小文 2017", "PDF 约 50% 位置")]),

    (20, "初阶方案 · Phase 0b", "Phase 0b：从已有论文提取的 10 组定量数据", "table",
     "10行数据表 — 四阶段温度节点/碳氧化窗口/活化窗口/N₂抗弯+43%/N₂导电率×10/N₂导热+0.62/脱羟基起始800°C/脱碳率99.93%/相转变T950 CI=27.3%/T1000介孔4.0nm — 来源PAPER-005/006/007",
     []),

    (21, "初阶方案 · Phase 0b", "Phase 0b：这篇论文的创新点在哪？", "two-col",
     "左：已有论文(PAPER-005仅空气/PAPER-006仅单温度) | 右：我们的贡献(温度×气氛=三位一体相图 · Data Re-analysis类型 · 不需新实验但产出新框架)",
     []),

    (22, "初阶方案 · Phase 0c", "Phase 0c：管式炉验证实验", "two-col",
     "左：实验矩阵SVG(4温度×3N₂流量=12组) + 设备照片占位 | 右：参数表(原料200目/保温2h/升温10°C/min/N₂气氛) · 预算~5,000元/2-3月",
     [("管式炉设备照片", "学校实验室实拍 或 网图", ""),
      ("黑滑石 200目粉体照片", "自拍", "")]),

    (23, "初阶方案 · Phase 0c", "Phase 0c：三个增量创新 — PAPER-006 没做过的", "cards",
     "三列卡片 — (1)Raman ID/IG定量石墨烯缺陷密度 vs 温度 vs 流量 / (2)N₂流量梯度50/100/200 mL/min / (3)高碳(~2%)vs 低碳(~0.5%)黑滑石对比确定经济边界",
     []),

    (24, "初阶方案 · Phase 0c", "Phase 0c：预算 ¥5,000 怎么分配", "two-col",
     "左：饼图SVG(管式炉40%/Raman24%/N₂16%/原料10%/XRD10%) | 右：分类说明表格 · 总计¥5,000",
     []),

    (25, "初阶方案 · Phase 0c", "Phase 0c：预期结果与论文规划", "two-col",
     "左：预期结果(ID/IG vs T曲线/ID/IG vs 流量/电阻率 vs T vs 流量/XRD物相演化) | 右：论文结构+目标期刊Ceramics International + 2个趋势图占位",
     [("ID/IG vs Temperature 预期趋势", "文献参考：典型碳材料Raman数据", "碳材料S曲线"),
      ("电阻率 vs Temperature 预期趋势", "文献参考：percolation曲线", "")]),

    (26, "初阶方案 · Phase 0d", "Phase 0d：发明专利申请", "two-col",
     "左：技术流程SVG(黑滑石→200目→N₂两阶段控温→N₂冷却→球磨分级→导电填料) | 右：三个创新点(矿石自身层间碳/四阶段精确控温/单步法)",
     []),

    (27, "初阶方案 · Phase 0d", "Phase 0d：现有专利格局 — 我们的位置", "table",
     "已有专利(CN1168795抗病毒/CN109851320A抗菌陶瓷/CN111517649B抗菌岩板/CN111467264A化妆品) + D3方向无专利覆盖→先发优势窗口",
     []),

    (28, "初阶方案 · 时间线", "Phase 0 总览：8 个月 · 4 项产出", "svg-chart",
     "甘特图 — Phase 0a(1-2月)→0b(2-3月)→0c(3-6月)→0d(5-7月) · 产出:综述初稿+相图论文+Raman实验论文+发明专利×1",
     []),

    (29, "初阶方案 · 路线图", "初阶路线图：从文献到企业验证的完整路径", "svg-chart",
     "四学期时间线 — 学期1-2:文献综述+数据再分析/学期3-4:管式炉实验+论文专利/学期5-6:回转窑N₂改造/学期7-8:D1煅烧优化 → 初阶成果检验 → 进入中阶",
     []),

    (30, "初阶方案 · 企业合作", "如何找到第一个合作企业", "list",
     "企业画像(本地煅烧厂) + 接触渠道(导师人脉⭐⭐⭐⭐⭐/矿区走访⭐⭐⭐⭐/科技局⭐⭐⭐⭐/行业协会⭐⭐⭐/1688⭐⭐⭐) + Phase 1企业投入~¥1-5万/次",
     [("广丰地区煅烧厂照片", "百度街景 或 实地拍摄", ""),
      ("广丰区地图 标注截图", "百度地图 标注截图", "")]),

    (31, "初阶方案 · 风险", "初阶风险矩阵：每件事的 Plan B", "table",
     "5项风险(学校无管式炉/N₂未石墨烯化/综述不够论文/专利被驳回/找不到企业) × 概率/影响/应对策略 · 核心原则:每一步都设计了不需要企业的备选路径",
     []),

    (32, "初阶方案 · 检验", "初阶完成标准：什么时候可以进入中阶？", "svg-chart",
     "四道门递进图 — 论文发表(1-2篇)→专利授权(1-2项)→企业意向(Phase 1反馈)→中试数据(公斤→百公斤) → 进入中阶合作谈判",
     []),

    (33, "初阶方案 · 身份", "初阶完成后的身份转换", "two-col",
     "左：现在(对黑滑石感兴趣的大一学生) → 右：初阶完成后(N₂气氛黑滑石碳石墨烯化方向有论文+专利的初级研究者) · 这个身份转换=进入中阶的门票",
     []),

    # === 中阶简览 (34-38) ===
    (34, "中阶方案 · 总引", "中阶方案：从课题组到产学研合作", "two-col",
     "初阶(已完成:论文+专利/¥8,000) → 中阶(本阶段:中试线/客户开发/500-5,000万) · 核心策略:D2(5G陶瓷)+D3(N₂碳复合)双旗舰并行",
     []),

    (35, "中阶方案 · 掺杂", "中阶核心决策：D2 路线选哪种掺杂？", "table",
     "5种掺杂对比(无掺杂Q×f=56,782/Cu 5% Q×f=93,600/Co 15% Q×f=145,846/Ge 4% Q×f=216,880) · 推荐:Cu(3%)+Co(8%)共掺杂 预期Q×f>120,000 |τf|<±15",
     [("Q×f vs τf 散点图", "PAPER-010 Fig.7 + PAPER-011 Fig.8", "PDF 后半部分")]),

    (36, "中阶方案 · 双旗舰", "中阶双旗舰：D3 先盈利，D2 建壁垒", "svg-chart",
     "双通道时间线 — D3快赢(吨级中试0-6月→送样6-12→订单12-18→扩产18-24月 CAPEX~50万 毛利率~93%) / D2高壁垒(配方定型0-6→中试6-18→基站验证18-30月 CAPEX~500万)",
     []),

    (37, "中阶方案 · 财务", "中阶经济模型概要", "table",
     "D3 vs D2对比表 — CAPEX(50-80万 vs 300-500万)/产能(100-500吨 vs 10-50吨)/单价(0.5-2万 vs 5-10万/吨)/年营收(50-500万 vs 50-500万)/毛利率(70-93% vs 40-60%) · 合计~500万/~500-1000万",
     []),

    (38, "中阶方案 · 触发", "中阶关键风险与进入条件", "list",
     "关键风险(Co供应链/客户验证周期/环保审批) + 四条触发信号(论文+专利+企业意向+中试数据) · 配套文件:tier-2-mid-level.md",
     []),

    # === 高阶简览 (39-41) ===
    (39, "高阶方案 · 愿景", "高阶方案：产学研联合体的全产业链整合", "list",
     "五条路线全线启动 — D1 50万吨/D2量产/D3出口/D4独立产品线/D9全线 · 总计:CAPEX~5.5亿/年营收~15亿/年EBIT~9.1亿/原料消化百万吨级",
     []),

    (40, "高阶方案 · 分级", "原料分级：同一座矿，五种产品", "svg-chart",
     "树状分配图 — 原矿100%→~5%特选(MgO>30%)→D4/~20%高镁(MgO28-31%)→D2/~10%高碳(>2%)→D3/~50%中等→D9/~15%低品位→D1 · 核心逻辑:什么矿适合什么产品",
     [("XRF 成分对比表", "PAPER-005 Table 1 + PAPER-007 Table 1", "PDF 前 15% 位置")]),

    (41, "高阶方案 · 协同", "五条路线的协同效应", "svg-chart",
     "废弃物互消化流程 — D1尾矿/D3筛下料/D9残渣→其他路线消化 · 关键协同:D3+D1共用煅烧窑/D2+D9原料降本/D9副产SiO₂→D1/D4 · 五线协同成本降~35-40% · 证据:PAPER-008成本降80%/PAPER-007 Mg97.98%+SiO₂99.92%",
     []),

    # === 收尾 (42-44) ===
    (42, "收尾 · 行动", "下一步行动：本周就能开始的三件事", "list",
     "立即(本周):确认设备/启动检索/下载Ge论文 · 短期(1-4周):综述提纲/采购样品 · 中期(1-6月):数据论文/管式炉实验/专利撰写",
     []),

    (43, "收尾 · 讨论", "希望得到导师指导的四个问题", "list",
     "四个结构化问题 — (1)学校是否有N₂气氛高温实验设备? / (2)是否有广丰本地煅烧企业联系人? / (3)D3(石墨烯填料)还是D2(5G陶瓷)更适合本科起点? / (4)是否可以安排研究生参与或作为本科毕设?",
     []),

    (44, "收尾 · 致谢", "感谢导师的宝贵时间", "dark",
     "致谢页 — 矿区远景照片背景 + 汇报人/[学校邮箱]/2026年7月 + 底部文字:10亿吨黑滑石,从一个本科生做起",
     [("矿区远景", "广丰矿区或江西丘陵地貌风景照", "")]),
]


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]

    for page_num, section_tag, title, slide_type, desc, placeholders in SLIDES:
        slide = prs.slides.add_slide(blank_layout)

        # 背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        if slide_type in ('cover', 'dark'):
            fill.fore_color.rgb = DEEP_BLUE
        else:
            fill.fore_color.rgb = TITANIUM

        # 章节标签 (左上)
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(4), Inches(0.35))
        tf = tag_box.text_frame
        tf.text = f"■ {section_tag}"
        p = tf.paragraphs[0]
        p.font.size = TAG_SIZE
        p.font.color.rgb = AMBER
        p.font.name = FONT_TITLE

        # 页面标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.text = title
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = TITLE_SIZE
        p.font.bold = True
        p.font.name = FONT_TITLE
        if slide_type in ('cover', 'dark'):
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        else:
            p.font.color.rgb = DEEP_BLUE

        # 内容描述
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True

        type_names = {
            'cover': '[封面页]', 'dark': '[暗底页]', 'two-col': '[双栏布局]',
            'cards': '[卡片矩阵]', 'table': '[数据表格]', 'svg-chart': '[SVG 图表]',
            'list': '[清单列表]'
        }

        p = tf.paragraphs[0]
        p.text = f"{type_names.get(slide_type, '')}  {desc}"
        p.font.size = Pt(14)
        p.font.name = FONT_BODY
        if slide_type in ('cover', 'dark'):
            p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        else:
            p.font.color.rgb = GRAPHITE

        # 占位框
        if placeholders:
            ph_top = Inches(3.0)
            for i, (ph_label, ph_source, ph_pos) in enumerate(placeholders):
                left_offset = Inches(0.8) if len(placeholders) < 3 else Inches(0.8 + (i % 2) * 5.8)
                top_offset = ph_top if len(placeholders) < 3 else ph_top + Inches((i // 2) * 1.3)

                ph_box = slide.shapes.add_shape(
                    1,  # MSO_SHAPE.RECTANGLE
                    left_offset, top_offset,
                    Inches(5.2), Inches(1.0)
                )
                ph_box.fill.solid()
                ph_box.fill.fore_color.rgb = RGBColor(0xF7, 0xF6, 0xF3)
                ph_box.line.color.rgb = LIGHT_GRAY
                ph_box.line.width = Pt(1.5)
                ph_box.line.dash_style = 2

                tf = ph_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = f"[图表占位] {ph_label}"
                p.font.size = Pt(12)
                p.font.name = FONT_TITLE
                p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)

                p2 = tf.add_paragraph()
                p2.text = f"来源: {ph_source}"
                p2.font.size = Pt(9)
                p2.font.name = FONT_BODY
                p2.font.color.rgb = RGBColor(0xB0, 0xB7, 0xC3)

                if ph_pos:
                    p3 = tf.add_paragraph()
                    p3.text = f"位置: {ph_pos}"
                    p3.font.size = Pt(9)
                    p3.font.name = FONT_BODY
                    p3.font.color.rgb = RGBColor(0xB0, 0xB7, 0xC3)

        # 页脚
        footer_top = Inches(7.0)
        footer_box = slide.shapes.add_textbox(Inches(8.5), footer_top, Inches(4.5), Inches(0.35))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{page_num} / 44"
        p.font.size = Pt(9)
        p.font.name = FONT_BODY
        p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
        p.alignment = PP_ALIGN.RIGHT

        src_box = slide.shapes.add_textbox(Inches(0.8), footer_top, Inches(7.5), Inches(0.35))
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = "数据来源: tier-1-entry-level.md + 对应论文分析报告"
        p.font.size = Pt(9)
        p.font.name = FONT_BODY
        p.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'advisor-presentation.pptx')
    prs.save(output_path)
    print(f"PPTX saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path


if __name__ == '__main__':
    create_presentation()
